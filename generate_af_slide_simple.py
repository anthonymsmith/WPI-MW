"""
Generate AF_Spring_Campaign_Slide_Simple.pptx — plain bullet-style committee briefing.
Two slides, white background, minimal formatting.
Each slide's body is a single text box with indented paragraphs for easy editing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

OUTPUT = '/Users/antho/Documents/WPI-MW/AF_Spring_Campaign_Slide_Simple.pptx'

NAVY  = RGBColor(0x1A, 0x3A, 0x5C)
TEAL  = RGBColor(0x2A, 0x9E, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY = RGBColor(0x33, 0x33, 0x33)
MGRAY = RGBColor(0x77, 0x77, 0x77)

# Indent levels in EMUs (1 inch = 914400)
IN0 = 0        # flush left
IN1 = 274320   # ~0.3"
IN2 = 548640   # ~0.6"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 0.9, NAVY)
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.05), Inches(12.6), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = False
    # Title
    p = tf.paragraphs[0]
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    # Subtitle on second paragraph in the same box
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(11)
        run2.font.italic = True
        run2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xDD)


def make_content_box(slide, left, top, width, height):
    """Create a single text box and return its text frame."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, size, bold, color, italic=False,
             space_before=0, space_after=3, left_emu=IN0):
    """Append a paragraph to tf with the given style and left-margin indent."""
    # First paragraph already exists; all subsequent ones need add_paragraph()
    if tf.paragraphs[0].runs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    if left_emu:
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(left_emu))
        pPr.set('indent', '0')
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


def gap(tf, pts=4):
    """Add a blank spacing paragraph."""
    add_para(tf, '', pts, False, DGRAY, space_before=0, space_after=0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Goal & Approach
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s1, 'Annual Fund — Spring Campaign',
          'Growing participation through targeted personal outreach')

tf1 = make_content_box(s1, 0.35, 1.05, 12.6, 5.9)

add_para(tf1, 'THE GOAL', 9, True, TEAL)
add_para(tf1,
    'Increase Annual Fund revenue and participation this spring by directing personal '
    'outreach to the patrons most likely to give — and making sure the right ask reaches the right person.',
    11, False, DGRAY, left_emu=IN1, space_after=8)

add_para(tf1, 'THE APPROACH', 9, True, TEAL, space_before=4)
add_para(tf1, '1.  Look at the data', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf1,
    'Four years of Annual Fund transaction history — who gave, how much, '
    'and whether they typically give in spring or fall.',
    10, False, DGRAY, left_emu=IN2, space_after=6)

add_para(tf1, '2.  Combine with patron behavior', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf1,
    'Attendance frequency, ticket spending, and loyalty patterns '
    'show how engaged each patron is with Music Worcester overall.',
    10, False, DGRAY, left_emu=IN2, space_after=6)

add_para(tf1, '3.  Segment for the right ask', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf1,
    'Group prospects by their giving relationship — renewal, reactivation, '
    'or first-time gift — so every outreach has a specific, relevant message.',
    10, False, DGRAY, left_emu=IN2, space_after=8)

add_para(tf1, 'RESULT', 9, True, TEAL, space_before=4)
add_para(tf1,
    'A focused list of ~200 priority contacts for personal outreach — ranked by giving history, '
    'engagement, and donation level — plus an extended list for email and letter follow-up.',
    11, False, DGRAY, left_emu=IN1)

# Footer
tb_f = s1.shapes.add_textbox(Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.25))
p_f = tb_f.text_frame.paragraphs[0]
p_f.alignment = PP_ALIGN.CENTER
run_f = p_f.add_run()
run_f.text = 'Data sources: Annual Fund transaction history FY2023–26  ·  Patron engagement data (attendance, spending, loyalty)'
run_f.font.size = Pt(8)
run_f.font.italic = True
run_f.font.color.rgb = MGRAY


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Outreach List
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
title_bar(s2, 'Your Outreach List',
          '~205 priority contacts for personal outreach  ·  ~280 for email/letter follow-up')

tf2 = make_content_box(s2, 0.35, 1.05, 12.6, 5.8)

add_para(tf2,
    'Start with these — highest-return outreach contacts, sorted by donor level within each group.',
    10, False, MGRAY, italic=True, space_after=8)

add_para(tf2, 'PERSONAL OUTREACH  (~205 contacts)', 9, True, TEAL, space_before=2)
add_para(tf2, '128   Spring Renew', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf2,
    'Gave last year; spring-aligned giving season. Renewal ask — highest conversion.',
    10, False, DGRAY, left_emu=IN2, space_after=5)

add_para(tf2, '54   Fall Renew  (Major & Beethoven)', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf2,
    'Gave last year; fall-leaning donors — but major donors often respond year-round.',
    10, False, DGRAY, left_emu=IN2, space_after=5)

add_para(tf2, '23   Reactivate  (Major & Beethoven)', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf2,
    'Gave $1,000+ in a prior year, now lapsed. A personal outreach matters most here.',
    10, False, DGRAY, left_emu=IN2, space_after=8)

add_para(tf2, 'EMAIL & LETTER FOLLOW-UP  (~280 contacts)', 9, True, TEAL, space_before=4)
add_para(tf2, '~197   Fall Renew  (Mid & Small)', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf2,
    'Remaining fall donors below Beethoven level.',
    10, False, DGRAY, left_emu=IN2, space_after=5)

add_para(tf2, '~81   Reactivate  (Mid & Small)', 11, True, NAVY, left_emu=IN1, space_after=1)
add_para(tf2,
    'Remaining lapsed donors — good for a letter campaign.',
    10, False, DGRAY, left_emu=IN2)

# Footer
tb_f2 = s2.shapes.add_textbox(Inches(0.35), Inches(7.0), Inches(8.5), Inches(0.32))
p_f2 = tb_f2.text_frame.paragraphs[0]
run_f2 = p_f2.add_run()
run_f2.text = ('Full detail in AF_Spring_Campaign.xlsx — Personal Outreach sheet. '
               'Each row includes email, giving history, region, and subscriber status.')
run_f2.font.size = Pt(8)
run_f2.font.italic = True
run_f2.font.color.rgb = MGRAY

tb_f3 = s2.shapes.add_textbox(Inches(9.2), Inches(7.0), Inches(3.8), Inches(0.32))
p_f3 = tb_f3.text_frame.paragraphs[0]
p_f3.alignment = PP_ALIGN.RIGHT
run_f3 = p_f3.add_run()
run_f3.text = 'Ticket-return donors (827) available as a separate acquisition list.'
run_f3.font.size = Pt(8)
run_f3.font.italic = True
run_f3.font.color.rgb = MGRAY

prs.save(OUTPUT)
print(f'✓ Saved: {OUTPUT}')
