"""
BACHtoberfest + Bach Birthday Bash YoY committee slides — plain-English version.

BACHtoberfest: stacked by slot (headliner / organ / choir / CONCORA),
since each year's lineup maps cleanly onto the same four roles.

Birthday Bash: stacked by just two categories — Headliner vs Supporting
programs — because the supporting lineup changed materially year-to-year
and slot-level comparison is misleading. The supporting event names are
listed in the slide bullets so readers can see what was in each year.
"""
import os
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from forecast_2526_comparison import load_data

os.chdir("/Users/antho/Documents/WPI-MW")

NAVY   = "#1A3A5C"
TEAL   = "#2A9EA0"
GOLD   = "#C4A35A"
ORANGE = "#E8922A"
SLATE  = "#8A9BAE"
LGRAY  = "#E8EDF2"
DGRAY  = "#4A5568"

RGB_NAVY  = RGBColor(0x1A, 0x3A, 0x5C)
RGB_TEAL  = RGBColor(0x2A, 0x9E, 0xA0)
RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RGB_DGRAY = RGBColor(0x33, 0x33, 0x33)
RGB_MGRAY = RGBColor(0x77, 0x77, 0x77)

OUT_PPTX = "/Users/antho/Documents/WPI-MW/TCB_YoY_Slides.pptx"
CHART_BFEST = "/Users/antho/Documents/WPI-MW/chart_bachtoberfest_yoy.png"
CHART_BBASH = "/Users/antho/Documents/WPI-MW/chart_birthday_bash_yoy.png"


# ── Data ──────────────────────────────────────────────────────────────────
def pull_series():
    em, merged = load_data()
    live = merged[
        (merged["EventType"] == "Live")
        & (merged["EventStatus"] == "Complete")
        & (merged["TicketStatus"] == "Active")
        & (merged["Quantity"] > 0)
    ].copy()

    def classify(name):
        n = str(name).lower()
        if "bachtoberfest" in n:
            return "BACHtoberfest"
        if "bach bash" in n or "birthday bash" in n:
            return "Birthday Bash"
        return None

    live["Series"] = live["EventName"].apply(classify)
    sub = live[live["Series"].notna()]
    return (sub.groupby(["Series", "Season", "EventName"], as_index=False)
               .agg(Tickets=("Quantity", "sum")))


# ── BACHtoberfest: slot-level stack (comparable year over year) ──────────
BFEST_SLOTS = [
    ("Headliner recital",   "BACHtoberfest 2024: Zlatomir Fung",
                             "BACHtoberfest 2025: Simone Dinnerstein Recital", NAVY),
    ("Organ concert",       "BACHtoberfest 2024: Bach Organ Multimedia Concert",
                             "BACHtoberfest 2025: Organ",                       TEAL),
    ("Choir concert",       "BACHtoberfest 2024 Choir",
                             "BACHtoberfest 2025: Choir",                       GOLD),
    ("CONCORA program",     "BACHtoberfest 2024: CONCORA",
                             "BACHtoberfest 2025: CONCORA, Baroklyn, Simone Dinnerstein", ORANGE),
]


# ── Birthday Bash: just "Headliner" + "Supporting programs" ──────────────
BBASH_HEADLINER = (
    "Bach Bash 2025: Jeremy Denk",
    "Bach's Birthday Bash 2026: The Sebastians",
)
BBASH_OTHER = {
    "2024-25": [
        "Bach Bash 2025: Handel + Haydn Society",
        "Bach Bash 2025: Worcester Chorus Secular Cantatas",
        "Bach Bash 2025: Worcester Bach Collective | Cantatathon",
    ],
    "2025-26": [
        "Bach's Birthday Bash 2026: Keyboards Up Close",
        "Bach's Birthday Bash 2026: Worcester Chamber Music Society",
        "Bach's Birthday Bash 2026: Cantatathon",
    ],
}


# ── Chart: stacked columns ────────────────────────────────────────────────
def stacked_values_bfest(agg):
    lookup = agg.set_index("EventName")["Tickets"].to_dict()
    rows = []
    for label, old_name, new_name, color in BFEST_SLOTS:
        rows.append({
            "Slot": label,
            "Old": int(lookup.get(old_name, 0)),
            "New": int(lookup.get(new_name, 0)),
            "Color": color,
        })
    return pd.DataFrame(rows)


def stacked_values_bbash(agg):
    lookup = agg.set_index("EventName")["Tickets"].to_dict()
    hl_old = int(lookup.get(BBASH_HEADLINER[0], 0))
    hl_new = int(lookup.get(BBASH_HEADLINER[1], 0))
    other_old = sum(int(lookup.get(n, 0)) for n in BBASH_OTHER["2024-25"])
    other_new = sum(int(lookup.get(n, 0)) for n in BBASH_OTHER["2025-26"])
    return pd.DataFrame([
        {"Slot": "Headliner night",     "Old": hl_old,    "New": hl_new,    "Color": NAVY},
        {"Slot": "Supporting programs", "Old": other_old, "New": other_new, "Color": SLATE},
    ])


def make_stacked(df, series_label, out_path, caption=None):
    fig, ax = plt.subplots(figsize=(7.2, 4.8))

    years = ["2024–25", "2025–26"]
    x = np.arange(2)
    width = 0.52

    bottom = np.zeros(2)
    for _, row in df.iterrows():
        vals = np.array([row["Old"], row["New"]])
        ax.bar(x, vals, width=width, bottom=bottom, color=row["Color"],
               edgecolor="white", linewidth=1.6, label=row["Slot"], zorder=3)
        for xi, v, b in zip(x, vals, bottom):
            if v >= 120:
                ax.text(xi, b + v / 2, f"{v:,}", ha="center", va="center",
                        fontsize=11, color="white", fontweight="bold")
        bottom += vals

    totals = [df["Old"].sum(), df["New"].sum()]
    ymax = max(totals) * 1.22
    ax.set_ylim(0, ymax)
    for xi, t in zip(x, totals):
        ax.text(xi, t + ymax * 0.02, f"{t:,}", ha="center", va="bottom",
                fontsize=14, fontweight="bold", color=NAVY)

    ax.set_xticks(x); ax.set_xticklabels(years, fontsize=12, color=DGRAY)
    ax.set_yticks([])
    for spine in ("top", "right", "left"):
        ax.spines[spine].set_visible(False)
    ax.spines["bottom"].set_color(LGRAY)

    delta_pct = (totals[1] - totals[0]) / totals[0] * 100
    ax.set_title(f"{series_label} — Total Attendance",
                 fontsize=14, fontweight="bold", color=NAVY,
                 loc="left", pad=28)
    # Subtitle anchored just above the plot area in axes coords.
    ax.text(0, 1.02,
            f"{totals[0]:,} tickets → {totals[1]:,} tickets  ({delta_pct:+.0f}%)",
            transform=ax.transAxes, fontsize=11, color=DGRAY,
            ha="left", va="bottom")

    ax.legend(loc="center left", bbox_to_anchor=(1.0, 0.5),
              fontsize=10.5, frameon=False, title="",
              handlelength=1.6, handleheight=1.6)

    if caption:
        fig.text(0.02, 0.005, caption, fontsize=8.5, color=DGRAY, style="italic")

    fig.tight_layout()
    fig.savefig(out_path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ── PPTX ──────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def title_bar(slide, title):
    add_rect(slide, 0, 0, 13.33, 0.95, RGB_NAVY)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = RGB_WHITE


def add_para(tf, text, size, bold, color, space_before=0, space_after=6):
    if tf.paragraphs[0].runs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", "0"); pPr.set("indent", "0")
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color


def build_slide(prs, title, bullets, chart_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, title)

    tb = s.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(5.9), Inches(5.9))
    tf = tb.text_frame; tf.word_wrap = True
    for label, body in bullets:
        if label:
            add_para(tf, label, 16, True, RGB_TEAL, space_before=4, space_after=2)
        add_para(tf, body, 18, False, RGB_DGRAY, space_after=12)

    s.shapes.add_picture(chart_path, Inches(6.55), Inches(1.2), width=Inches(6.5))


def build_text_slide(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, title)

    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(6.1))
    tf = tb.text_frame; tf.word_wrap = True
    for label, body in bullets:
        if label:
            add_para(tf, label, 18, True, RGB_TEAL, space_before=6, space_after=3)
        add_para(tf, body, 18, False, RGB_DGRAY, space_after=14)


def main():
    agg = pull_series()
    bf = stacked_values_bfest(agg)
    bb = stacked_values_bbash(agg)

    print("BACHtoberfest:")
    print(bf.to_string(index=False))
    print("\nBirthday Bash:")
    print(bb.to_string(index=False))

    make_stacked(bf, "BACHtoberfest", CHART_BFEST)
    make_stacked(bb, "Bach Birthday Bash", CHART_BBASH,
                 caption="Supporting programs differed between years — see slide for the lineups.")
    print(f"\n✓ {CHART_BFEST}\n✓ {CHART_BBASH}")

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    bf_bullets = [
        ("THE HEADLINE",
         "About 300 fewer tickets sold across BACHtoberfest this year "
         "(1,823 → 1,520)."),
        ("WHAT GREW",
         "The CONCORA program drew a larger audience after the slot expanded "
         "to include Baroklyn and Simone Dinnerstein."),
        ("WHAT SLIPPED",
         "The headliner recital was the biggest single drop. The organ and "
         "choir concerts each brought in a little less than last year."),
        ("TAKEAWAY",
         "The CONCORA bundle with Baroklyn and Dinnerstein grew the slot — "
         "proof that repackaging works. But Dinnerstein also headlined the "
         "solo recital, and that night took the biggest single drop. Looks "
         "like marquee-name stacking cannibalized the recital."),
    ]
    build_slide(prs, "BACHtoberfest — Last Year vs This Year", bf_bullets, CHART_BFEST)

    bb_bullets = [
        ("THE HEADLINE",
         "About 240 fewer tickets sold across Bach Birthday Bash this year "
         "(1,772 → 1,531)."),
        ("HEADLINER GREW",
         "The Sebastians drew a bigger house than last year's Jeremy Denk "
         "— the marquee night is the strongest it has been."),
        ("LINEUP CHANGED",
         "The three supporting programs were different events this year. "
         "Last year: Handel + Haydn Society, Worcester Chorus Secular "
         "Cantatas, Cantatathon. This year: Keyboards Up Close, Worcester "
         "Chamber Music Society, Cantatathon."),
        ("TAKEAWAY",
         "The strong marquee (Sebastians) held up. The three new supporting "
         "programs are effectively year one — judge them on a two-to-three "
         "year build, not a year-one comparison against last year's lineup."),
    ]
    build_slide(prs, "Bach Birthday Bash — Last Year vs This Year", bb_bullets, CHART_BBASH)

    sowhat_bullets = [
        ("THE CONSTRAINT",
         "TCB has to work through the full Bach repertoire — a lot of "
         "cantatas and organ works alongside a much smaller number of "
         "orchestral and solo-piano headliners. The bottleneck isn't "
         "audience demand; it's the supply of high-draw slots."),
        ("1. SPACE THE MARQUEES ACROSS FESTIVALS",
         "Distribute headliner-tier names across the whole TCB calendar, "
         "not within a single festival. BACHtoberfest 2025 used Dinnerstein "
         "in two slots (solo recital + CONCORA bundle); Birthday Bash had "
         "one marquee. Every festival should carry at least one anchor; no "
         "festival should double-book the same name."),
        ("2. BUNDLE THE MISSION SLOTS",
         "The CONCORA result is the template. Standalone cantata and organ "
         "nights are where the softness lives. Frame them around an anchor "
         "artist or a theme — the way CONCORA was repackaged with Baroklyn "
         "and Dinnerstein — rather than letting them stand alone."),
        ("3. ANCHOR-PLUS-HALO PRICING",
         "Bundle cantata and organ tickets with the headliner, or use "
         "pay-what-you-will on the mission slots. Historical free/PWYW "
         "programs have drawn roughly 2.9× their paid-pricing baseline — "
         "a real lever for getting Bach-curious audiences in the room."),
    ]
    build_text_slide(prs, "So What — Recommendations", sowhat_bullets)

    prs.save(OUT_PPTX)
    print(f"✓ {OUT_PPTX}")


if __name__ == "__main__":
    main()
