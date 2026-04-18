"""
Generate AF_Spring_Campaign_Slide.pptx — 2-slide revenue committee briefing.
  Slide 1: Goal & approach
  Slide 2: The prioritized outreach list (framed as manageable)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = '/Users/antho/Documents/WPI-MW/AF_Spring_Campaign_Slide.pptx'

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
TEAL   = RGBColor(0x2A, 0x9E, 0xA0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
MGRAY  = RGBColor(0x88, 0x88, 0x88)
GREEN  = RGBColor(0xD5, 0xE8, 0xD4)
BLUE   = RGBColor(0xDA, 0xE8, 0xFC)
AMBER  = RGBColor(0xFF, 0xE6, 0xCC)
GREEN_D  = RGBColor(0x38, 0x76, 0x1D)
BLUE_D   = RGBColor(0x1F, 0x51, 0x7A)
AMBER_D  = RGBColor(0xA0, 0x52, 0x00)
GOLD   = RGBColor(0xE8, 0x92, 0x2A)

def rgb(r,g,b): return RGBColor(r,g,b)

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill, line=None, line_w=0.5):
    from pptx.util import Pt
    s = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def textbox(slide, left, top, width, height, text,
            size=11, bold=False, italic=False,
            color=DGRAY, bg=None, align=PP_ALIGN.LEFT, wrap=True, spacing=None):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        from pptx.util import Pt as PPt
        p.space_after = PPt(spacing)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
    return tb

def title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    textbox(slide, 0.35, 0.1, 11, 0.55, title,
            size=26, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, 0.35, 0.65, 11, 0.35, subtitle,
                size=11, italic=True, color=rgb(0xAA,0xCC,0xDD))

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Goal & Approach
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])

title_bar(s1,
    'Annual Fund — Spring Campaign',
    'Growing participation through targeted personal outreach')

# Goal section
add_rect(s1, 0.35, 1.25, 12.6, 0.01, TEAL)
textbox(s1, 0.35, 1.32, 12.6, 0.32,
        'THE GOAL', size=9, bold=True, color=TEAL)
textbox(s1, 0.35, 1.60, 12.6, 0.55,
        'Increase Annual Fund revenue and participation this spring by directing personal outreach '
        'to the patrons most likely to give — and making sure the right ask reaches the right person.',
        size=13, color=NAVY, bold=False)

# Divider
add_rect(s1, 0.35, 2.28, 12.6, 0.01, LGRAY)

# Approach section
textbox(s1, 0.35, 2.36, 12.6, 0.32,
        'THE APPROACH', size=9, bold=True, color=TEAL)

# Three step boxes
steps = [
    ('1', 'Look at the data',
     'Four years of Annual Fund transaction history — who gave, how much, '
     'and whether they typically give in spring or fall.'),
    ('2', 'Combine with patron behavior',
     'Attendance frequency, ticket spending, and loyalty patterns '
     'show how engaged each patron is with Music Worcester overall.'),
    ('3', 'Segment for the right ask',
     'Group prospects by their giving relationship — renewal, reactivation, '
     'or first-time gift — so every outreach has a specific, relevant message.'),
]

box_w  = 3.85
box_h  = 2.2
box_y  = 2.75
gap    = 0.17
for i, (num, heading, body) in enumerate(steps):
    x = 0.35 + i * (box_w + gap)
    add_rect(s1, x, box_y, box_w, box_h, LGRAY, line=rgb(0xCC,0xCC,0xCC))
    # Big step number
    add_rect(s1, x, box_y, box_w, 0.52, TEAL if i==0 else NAVY if i==1 else GOLD)
    textbox(s1, x+0.15, box_y+0.04, 0.38, 0.42,
            num, size=22, bold=True, color=WHITE)
    textbox(s1, x+0.55, box_y+0.1, box_w-0.65, 0.34,
            heading, size=12, bold=True, color=WHITE)
    textbox(s1, x+0.18, box_y+0.62, box_w-0.3, 1.45,
            body, size=11, color=DGRAY)

# Result callout
add_rect(s1, 0.35, 5.10, 12.6, 0.85, rgb(0xE8,0xF4,0xF8), line=TEAL, line_w=1.5)
textbox(s1, 0.55, 5.18, 1.5, 0.65,
        'Result:', size=13, bold=True, color=TEAL)
textbox(s1, 1.85, 5.18, 10.9, 0.65,
        'A focused list of ~200 priority contacts for personal outreach — '
        'ranked by giving history, engagement, and donation level — '
        'plus an extended list for email and letter follow-up.',
        size=12, color=NAVY)

# Footer
textbox(s1, 0.35, 6.15, 12.6, 0.28,
        'Data sources: Annual Fund transaction history FY2023–26  ·  Patron engagement data (attendance, spending, loyalty)',
        size=8, italic=True, color=MGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Outreach List
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])

title_bar(s2,
    'Your Outreach List',
    '~205 priority contacts for personal calls  ·  ~280 for email/letter follow-up')

# Intro line
textbox(s2, 0.35, 1.2, 12.6, 0.38,
        'Start here — these three groups are your highest-return personal calls. '
        'Each is sorted within the spreadsheet: Major donors ($2,500+) and Beethoven level ($1K–$2.4K) first.',
        size=10, italic=True, color=DGRAY)

# ── Priority personal outreach (3 groups) ──
add_rect(s2, 0.35, 1.65, 12.6, 0.32, NAVY)
textbox(s2, 0.5, 1.70, 12.2, 0.25,
        'PRIORITY — Personal Outreach  (~205 contacts)',
        size=10, bold=True, color=WHITE)

card_data = [
    ('128', 'Spring Renew',
     'Gave to Annual Fund last year and typically give in spring',
     'Renewal ask — in-season, highest conversion',
     'All 128  ·  27 Major  ·  19 Beethoven', GREEN, GREEN_D),
    ('54', 'Fall Renew\n(Major & Beethoven)',
     'Gave last year; typically give in fall — but significant donors worth asking now',
     'Off-cycle ask — major donors often respond year-round',
     '54 of 251 fall donors  ·  29 Major  ·  25 Beethoven', BLUE, BLUE_D),
    ('23', 'Reactivate\n(Major & Beethoven)',
     'Gave $1,000+ in a prior year but not recently — lapsed but high-value',
     'Re-engagement — a personal call matters most here',
     '23 of 104 lapsed donors  ·  5 Major  ·  18 Beethoven', AMBER, AMBER_D),
]

card_w = 3.88
card_y = 2.05
card_h = 2.65
for i, (count, name, who, ask, detail, bg, fg) in enumerate(card_data):
    x = 0.35 + i * (card_w + 0.17)
    add_rect(s2, x, card_y, card_w, card_h, bg, line=rgb(0xCC,0xCC,0xCC))
    # Count badge
    add_rect(s2, x, card_y, 0.72, card_h, fg)
    textbox(s2, x+0.03, card_y+0.62, 0.65, 0.55,
            count, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s2, x+0.03, card_y+1.18, 0.65, 0.35,
            'contacts', size=7, color=WHITE, align=PP_ALIGN.CENTER)
    # Content
    textbox(s2, x+0.8, card_y+0.1, card_w-0.9, 0.42,
            name, size=11, bold=True, color=NAVY)
    textbox(s2, x+0.8, card_y+0.52, card_w-0.9, 0.65,
            who, size=9, color=DGRAY, italic=True)
    textbox(s2, x+0.8, card_y+1.18, card_w-0.9, 0.52,
            ask, size=10, bold=True, color=fg)
    add_rect(s2, x+0.78, card_y+1.74, card_w-0.88, 0.01, rgb(0xCC,0xCC,0xCC))
    textbox(s2, x+0.8, card_y+1.82, card_w-0.9, 0.72,
            detail, size=8, color=MGRAY)

# ── Email / letter follow-up ──
add_rect(s2, 0.35, 4.82, 12.6, 0.30, LGRAY, line=rgb(0xCC,0xCC,0xCC))
textbox(s2, 0.5, 4.87, 4.0, 0.23,
        'EXTENDED LIST — Email / Letter Follow-up  (~280 contacts)',
        size=9, bold=True, color=DGRAY)

ext_items = [
    ('~197', 'Fall Renew (Mid & Small)', 'Remaining fall donors below Beethoven level'),
    ('~81',  'Reactivate (Mid & Small)', 'Remaining lapsed donors — good for a letter campaign'),
]
for i, (cnt, label, note) in enumerate(ext_items):
    x = 0.35 + i * 6.18
    textbox(s2, x+0.15, 5.18, 0.55, 0.42,
            cnt, size=16, bold=True, color=TEAL)
    textbox(s2, x+0.72, 5.18, 2.5, 0.22,
            label, size=9, bold=True, color=DGRAY)
    textbox(s2, x+0.72, 5.40, 2.5, 0.22,
            note, size=8, italic=True, color=MGRAY)

# Footer note
textbox(s2, 0.35, 6.58, 8.5, 0.32,
        'Full detail in AF_Spring_Campaign.xlsx — Priority Outreach sheet. '
        'Each row includes email, giving history, region, and subscriber status.',
        size=8, italic=True, color=MGRAY)
textbox(s2, 9.2, 6.58, 3.8, 0.32,
        'Ticket-return donors (827) available as a separate acquisition list.',
        size=8, italic=True, color=MGRAY, align=PP_ALIGN.RIGHT)

prs.save(OUTPUT)
print(f'✓ Saved: {OUTPUT}')
