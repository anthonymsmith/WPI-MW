"""
Forecasting deck — tells the story of the attendance forecast model
for Music Worcester. Plain-English, non-technical committee style.

Five slides:
1. Title
2. The Case for Forecasting — why + what the model does
3. How It's Performing — 25-26 wide chart
4. What It's Unlocking — levers + decisions + compounding loop
5. What's Next — roadmap + closing
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

os.chdir("/Users/antho/Documents/WPI-MW")

RGB_NAVY   = RGBColor(0x1A, 0x3A, 0x5C)
RGB_TEAL   = RGBColor(0x2A, 0x9E, 0xA0)
RGB_TEAL_L = RGBColor(0x7A, 0xB8, 0xBA)
RGB_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RGB_DGRAY  = RGBColor(0x33, 0x33, 0x33)
RGB_MGRAY  = RGBColor(0x77, 0x77, 0x77)

OUT_PPTX    = "/Users/antho/Documents/WPI-MW/Forecasting_Deck.pptx"
CHART_25_26 = "/Users/antho/Documents/WPI-MW/forecast_2526_bar_chart_wide.png"


def add_rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def title_bar(slide, title):
    add_rect(slide, 0, 0, 13.33, 0.95, RGB_NAVY)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = RGB_WHITE


def add_para(tf, text, size, bold, color, space_before=0, space_after=6):
    if tf.paragraphs[0].runs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", "0"); pPr.set("indent", "0")
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color


def build_title_slide(prs, eyebrow, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 13.33, 7.5, RGB_NAVY)
    add_rect(s, 0.8, 5.45, 4.0, 0.05, RGB_TEAL)

    tb0 = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.5))
    p0 = tb0.text_frame.paragraphs[0]
    r0 = p0.add_run(); r0.text = eyebrow
    r0.font.size = Pt(14); r0.font.bold = True; r0.font.color.rgb = RGB_TEAL_L

    tb1 = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.2))
    p1 = tb1.text_frame.paragraphs[0]
    r1 = p1.add_run(); r1.text = title
    r1.font.size = Pt(48); r1.font.bold = True; r1.font.color.rgb = RGB_WHITE

    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(20); r2.font.color.rgb = RGB_TEAL_L


def build_text_slide(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, title)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(6.1))
    tf = tb.text_frame; tf.word_wrap = True
    for label, body in bullets:
        if label:
            add_para(tf, label, 18, True, RGB_TEAL, space_before=6, space_after=3)
        add_para(tf, body, 18, False, RGB_DGRAY, space_after=14)


def build_chart_slide(prs, title, subtitle, chart_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, title)

    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    add_para(tf, subtitle, 15, True, RGB_DGRAY, space_after=2)

    # Wide chart: aspect ≈ 2.04. Fit to slide width with side margins.
    pic_w = Inches(12.3)
    pic = s.shapes.add_picture(chart_path, Inches(0.5), Inches(1.85), width=pic_w)
    # Vertically center within remaining 5.6" of space
    avail_top = Inches(1.85)
    avail_h = Inches(7.5) - avail_top
    if pic.height < avail_h:
        pic.top = int(avail_top + (avail_h - pic.height) / 2)


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Title
    build_title_slide(
        prs,
        "MUSIC WORCESTER  ·  ATTENDANCE ANALYTICS",
        "Looking Ahead",
        "An attendance forecast built from your own history — and what to do with it.",
    )

    # 2. The Case — why + what combined
    case_bullets = [
        ("THE QUESTION",
         "How many people will show up?"),
        ("THE STAKES",
         "Every booking is a financial commitment — venue, artist, "
         "staffing, marketing — made months before a ticket goes on "
         "sale. Until now, those decisions have leaned on last year's "
         "numbers plus staff instinct."),
        ("THE APPROACH",
         "The model reads years of your own ticketing history to learn "
         "the shape of your audience — by venue, genre, event type, and "
         "artist profile. It adjusts for what's different about each "
         "booking, and is tested only on seasons it hasn't seen before."),
    ]
    build_text_slide(prs, "The Case for Forecasting", case_bullets)

    # 3. Results chart
    build_chart_slide(
        prs,
        "How It's Performing — 25–26 Season",
        "Across 26 completed events: predictions within about a quarter of actual on "
        "average; the season total lands within 1% of forecast.",
        CHART_25_26,
    )

    # 4. What It's Unlocking — levers + decisions + loop
    unlock_bullets = [
        ("QUANTIFIES THE LEVERS",
         "Pay-what-you-will roughly triples the paid-pricing baseline. "
         "Bundled programs (CONCORA + Baroklyn + Dinnerstein) grew the "
         "slot. Repeat series carry measurable loyalty. These used to "
         "be intuition — now they're numbers."),
        ("INFORMS THE DECISIONS",
         "Build the season budget from expected attendance per event. "
         "Direct marketing spend to events pacing soft, not the ones "
         "already on track. Make programming calls — bundle vs. "
         "standalone, marquee placement, pricing — with evidence "
         "behind them."),
        ("COMPOUNDS OVER TIME",
         "Every season's forecast-vs-actual comparison sharpens the "
         "next. The model is an asset, not a one-time report."),
    ]
    build_text_slide(prs, "What It's Unlocking", unlock_bullets)

    # 5. What's Next
    next_bullets = [
        ("PRICING OPTIMIZATION",
         "The next layer: sales over time. By watching how tickets sell "
         "in the weeks leading up to an event, the same framework can "
         "inform when and whether to raise, hold, or discount — and "
         "which tiers to offer."),
        ("SHARPER EVERY SEASON",
         "Each completed season adds training data. Artist signals "
         "grow. Repeat series accumulate history. Accuracy improves as "
         "the model sees more of your world."),
        ("ALREADY IN USE",
         "The 25–26 numbers you've seen this season came from this "
         "model. The TCB YoY story was built on the same foundation. "
         "It's not a future deliverable — it's running now."),
    ]
    build_text_slide(prs, "What's Next", next_bullets)

    prs.save(OUT_PPTX)
    print(f"✓ {OUT_PPTX}")


if __name__ == "__main__":
    main()
